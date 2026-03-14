#!/usr/bin/env python3
"""
Generate Chapter 1 PowerPoint presentation
Based on the specification in output/ppt/generate_ppt.md
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import os

# Color scheme
PURPLE_DARK = RGBColor(102, 126, 234)  # #667eea
PURPLE_DARKER = RGBColor(118, 75, 162)  # #764ba2
TEXT_DARK = RGBColor(51, 51, 51)  # #333333
WHITE = RGBColor(255, 255, 255)  # #ffffff

def create_presentation():
    """Create the PowerPoint presentation for Chapter 1"""

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Define slide layouts
    blank_layout = prs.slide_layouts[6]  # Blank layout

    slides_data = [
        {
            "title": "1️⃣ Git 시작하기",
            "subtitle": "Git의 기초부터 첫 저장소까지",
            "type": "title"
        },
        {
            "title": "Git이란 무엇인가?",
            "content": "• Git = 버전 관리 시스템\n• 파일의 변경 이력을 기록하고 관리\n\n\"Word 파일을 v1, v2...로 저장하는 대신\nGit이 효율적으로 관리해줍니다!\"",
            "type": "content"
        },
        {
            "title": "Git의 주요 장점",
            "content": "• 변경 기록 관리 - 누가, 언제, 무엇을 바꿨는지 기록\n• 이전 버전 복구 - 실수한 코드 언제든 복구\n• 협업 용이 - 여러 사람이 효율적으로 작업\n• 브랜치 기능 - 여러 작업을 병렬로 진행",
            "type": "content"
        },
        {
            "title": "Git의 구조",
            "content": "로컬 저장소\n↔️\n원격 저장소\n\n커밋 = 파일의 변경사항을 저장하는 행위",
            "type": "content"
        },
        {
            "title": "설치 - Windows",
            "content": "1️⃣ git-scm.com 방문\n\n2️⃣ \"Download for Windows\" 클릭\n\n3️⃣ 설치 파일 실행\n\n4️⃣ 기본 설정으로 설치 완료",
            "type": "content"
        },
        {
            "title": "설치 - macOS",
            "content": "터미널에서 다음 명령어 실행:\n\nbrew install git\n\n명령어 입력만 하면 됨!",
            "type": "content"
        },
        {
            "title": "설치 - Linux",
            "content": "Ubuntu/Debian:\nsudo apt-get install git\n\nFedora:\nsudo yum install git",
            "type": "content"
        },
        {
            "title": "설치 확인하기",
            "content": "다음 명령어로 확인:\n\ngit --version\n\n✅ \"git version X.X.X\" 메시지가\n나타나면 설치 완료!",
            "type": "content"
        },
        {
            "title": "사용자 정보 설정",
            "content": "Step 1: 사용자 이름 설정\ngit config --global user.name \"이름\"\n\nStep 2: 이메일 설정\ngit config --global user.email \"이메일\"",
            "type": "content"
        },
        {
            "title": "설정 확인하기",
            "content": "전체 설정 확인:\ngit config --list\n\n개별 확인:\ngit config user.name\ngit config user.email",
            "type": "content"
        },
        {
            "title": "첫 저장소 만들기",
            "content": "Step 1: 폴더 생성\nmkdir my-first-project\n\nStep 2: 폴더로 이동\ncd my-first-project",
            "type": "content"
        },
        {
            "title": "Git 저장소 초기화",
            "content": "Step 3: Git 저장소 초기화\n\ngit init\n\n.git 폴더가 생성됩니다.",
            "type": "content"
        },
        {
            "title": "저장소 상태 확인",
            "content": "git status\n\n출력 메시지:\nOn branch master\nNo commits yet\nnothing to commit",
            "type": "content"
        },
        {
            "title": "축하합니다! 🎉",
            "content": "첫 Git 저장소를 만들었습니다!\n\n✅ 완성!\n\n다음: 2장 기본 작업 (add, commit, push)",
            "type": "content"
        },
        {
            "title": "다음 단계",
            "content": "✨ 이제 할 수 있는 것:\n• 파일을 추가하고 변경사항 기록\n• 팀과 협업\n• 버전 관리\n\n🎓 2장에서 기본 명령어를 배워봅시다!",
            "type": "content"
        }
    ]

    # Create slides
    for i, slide_data in enumerate(slides_data, 1):
        if slide_data["type"] == "title":
            slide = create_title_slide(prs, blank_layout, slide_data)
        else:
            slide = create_content_slide(prs, blank_layout, slide_data)

        # Add page number
        add_page_number(slide, i)

    # Save the presentation
    output_dir = "output/ppt"
    os.makedirs(output_dir, exist_ok=True)
    output_file = os.path.join(output_dir, "Chapter1_Git_시작하기.pptx")
    prs.save(output_file)
    print(f"✅ PPT 생성 완료: {output_file}")
    print(f"   총 {len(slides_data)}개의 슬라이드")

def create_title_slide(prs, layout, data):
    """Create a title slide"""
    slide = prs.slides.add_slide(layout)

    # Add gradient background (approximated with colored shape)
    background = slide.shapes.add_shape(
        1,  # Rectangle
        0, 0,
        prs.slide_width,
        prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = PURPLE_DARK
    background.line.color.rgb = PURPLE_DARK

    # Add title
    title_box = slide.shapes.add_textbox(
        Inches(0.5),
        Inches(2.5),
        Inches(9),
        Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_p = title_frame.paragraphs[0]
    title_p.text = data["title"]
    title_p.font.size = Pt(54)
    title_p.font.bold = True
    title_p.font.color.rgb = WHITE
    title_p.alignment = PP_ALIGN.CENTER

    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(
        Inches(0.5),
        Inches(4.2),
        Inches(9),
        Inches(1)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    subtitle_p = subtitle_frame.paragraphs[0]
    subtitle_p.text = data["subtitle"]
    subtitle_p.font.size = Pt(28)
    subtitle_p.font.color.rgb = WHITE
    subtitle_p.alignment = PP_ALIGN.CENTER

    return slide

def create_content_slide(prs, layout, data):
    """Create a content slide with title and bullet points"""
    slide = prs.slides.add_slide(layout)

    # Add white background
    background = slide.shapes.add_shape(
        1,  # Rectangle
        0, 0,
        prs.slide_width,
        prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = WHITE
    background.line.color.rgb = WHITE

    # Add title
    title_box = slide.shapes.add_textbox(
        Inches(0.5),
        Inches(0.4),
        Inches(9),
        Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_p = title_frame.paragraphs[0]
    title_p.text = data["title"]
    title_p.font.size = Pt(40)
    title_p.font.bold = True
    title_p.font.color.rgb = PURPLE_DARK

    # Add horizontal line under title
    line = slide.shapes.add_connector(1, Inches(0.5), Inches(1.3), Inches(9.5), Inches(1.3))
    line.line.color.rgb = PURPLE_DARK
    line.line.width = Pt(2)

    # Add content
    content_box = slide.shapes.add_textbox(
        Inches(0.7),
        Inches(1.7),
        Inches(8.6),
        Inches(5.3)
    )
    content_frame = content_box.text_frame
    content_frame.word_wrap = True

    # Clear default paragraph
    content_frame.clear()

    # Add content text
    lines = data["content"].split("\n")
    for i, line in enumerate(lines):
        if i > 0:
            p = content_frame.add_paragraph()
        else:
            p = content_frame.paragraphs[0]

        p.text = line
        p.font.size = Pt(22) if line.startswith("•") else Pt(18)
        p.font.color.rgb = TEXT_DARK
        p.level = 0

        # Add spacing
        p.space_before = Pt(6)
        p.space_after = Pt(6)

    return slide

def add_page_number(slide, page_num):
    """Add page number to slide"""
    page_box = slide.shapes.add_textbox(
        Inches(9.2),
        Inches(7.1),
        Inches(0.7),
        Inches(0.3)
    )
    page_frame = page_box.text_frame
    page_p = page_frame.paragraphs[0]
    page_p.text = str(page_num)
    page_p.font.size = Pt(10)
    page_p.font.color.rgb = RGBColor(150, 150, 150)
    page_p.alignment = PP_ALIGN.RIGHT

if __name__ == "__main__":
    create_presentation()
