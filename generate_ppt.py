#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Git 입문 교재 PPT 생성 스크립트
HTML 원본을 읽어 PowerPoint 프레젠테이션을 생성합니다.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from html.parser import HTMLParser
import re
import os

# 색상 정의 (각 장별 고유 색상)
CHAPTER_COLORS = {
    1: {"primary": RGBColor(102, 126, 234), "secondary": RGBColor(118, 75, 162)},   # 보라
    2: {"primary": RGBColor(33, 150, 243), "secondary": RGBColor(25, 118, 210)},    # 파랑
    3: {"primary": RGBColor(255, 152, 0), "secondary": RGBColor(245, 124, 0)},      # 주황
    4: {"primary": RGBColor(244, 67, 54), "secondary": RGBColor(211, 47, 47)},      # 빨강
    5: {"primary": RGBColor(156, 39, 176), "secondary": RGBColor(106, 27, 154)},    # 자주
}

class HTMLTextExtractor(HTMLParser):
    """HTML에서 텍스트와 구조 추출"""

    def __init__(self):
        super().__init__()
        self.text_content = []
        self.in_code = False
        self.in_section = False
        self.sections = []
        self.current_section = {"title": "", "content": [], "code": []}

    def handle_starttag(self, tag, attrs):
        if tag == "section":
            self.in_section = True
        elif tag == "code" or tag == "pre":
            self.in_code = True

    def handle_endtag(self, tag):
        if tag == "section":
            if self.current_section["title"]:
                self.sections.append(self.current_section)
            self.current_section = {"title": "", "content": [], "code": []}
            self.in_section = False
        elif tag == "code" or tag == "pre":
            self.in_code = False

    def handle_data(self, data):
        text = data.strip()
        if text:
            if self.in_code:
                self.current_section["code"].append(text)
            elif self.in_section:
                if not self.current_section["title"] and text:
                    self.current_section["title"] = text
                else:
                    self.current_section["content"].append(text)

def extract_sections_from_html(html_file):
    """HTML 파일에서 섹션 추출"""
    with open(html_file, 'r', encoding='utf-8') as f:
        content = f.read()

    # h2 태그를 섹션 제목으로 인식
    sections = re.findall(r'<h2[^>]*>([^<]+)</h2>.*?(?=<h2|<footer|$)', content, re.DOTALL)

    # 더 정확한 섹션 추출: <section> 태그 이용
    section_blocks = re.findall(r'<section[^>]*>(.*?)</section>', content, re.DOTALL)

    extracted = []
    for block in section_blocks:
        # h2 제목 추출
        title_match = re.search(r'<h2[^>]*>([^<]+)</h2>', block)
        title = title_match.group(1) if title_match else "제목 없음"

        # p 태그 내용 추출 (최대 5개)
        paragraphs = re.findall(r'<p[^>]*>([^<]+)</p>', block)[:5]

        # code 블록 추출
        codes = re.findall(r'<code[^>]*>([^<]+)</code>', block)

        if title and paragraphs:
            extracted.append({
                "title": title.strip(),
                "content": [p.strip() for p in paragraphs if p.strip()],
                "code": [c.strip() for c in codes if c.strip()]
            })

    return extracted

def create_title_slide(prs, chapter_num, chapter_title, chapter_colors):
    """제목 슬라이드 생성"""
    blank_layout = prs.slide_layouts[6]  # 빈 레이아웃
    slide = prs.slides.add_slide(blank_layout)

    # 배경색
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = chapter_colors["primary"]

    # 장 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True

    p = title_frame.paragraphs[0]
    p.text = f"Chapter {chapter_num}"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # 부제목
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    p = subtitle_frame.paragraphs[0]
    p.text = chapter_title
    p.font.size = Pt(32)
    p.font.color.rgb = RGBColor(255, 255, 255)

    return slide

def create_content_slide(prs, section, chapter_colors, slide_num):
    """내용 슬라이드 생성"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # 흰색 배경
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

    # 상단 색상 바
    color_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.3))
    color_bar.fill.solid()
    color_bar.fill.fore_color.rgb = chapter_colors["primary"]
    color_bar.line.color.rgb = chapter_colors["primary"]

    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = section["title"]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(51, 51, 51)

    # 내용
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.6), Inches(4))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True

    for para_text in section["content"][:3]:  # 최대 3개 단락
        p = content_frame.add_paragraph()
        p.text = para_text[:100]  # 최대 100자
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(51, 51, 51)
        p.space_before = Pt(6)
        p.level = 0

    # 코드 블록
    if section["code"]:
        code_box = slide.shapes.add_textbox(Inches(0.7), Inches(5.5), Inches(8.6), Inches(1.5))
        code_frame = code_box.text_frame
        code_frame.word_wrap = True

        p = code_frame.paragraphs[0]
        p.text = "$ " + section["code"][0][:70]
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.name = 'Courier New'

        # 코드 배경
        code_shape = slide.shapes.add_shape(1, Inches(0.5), Inches(5.3), Inches(9), Inches(1.8))
        code_shape.fill.solid()
        code_shape.fill.fore_color.rgb = RGBColor(40, 40, 40)
        code_shape.line.color.rgb = RGBColor(150, 150, 150)

        # 코드 텍스트를 다시 앞에 놓기
        slide.shapes._spTree.remove(code_shape._element)
        slide.shapes._spTree.insert(2, code_shape._element)

    # 슬라이드 번호
    page_num_box = slide.shapes.add_textbox(Inches(9.2), Inches(6.8), Inches(0.7), Inches(0.3))
    p = page_num_box.text_frame.paragraphs[0]
    p.text = str(slide_num)
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(150, 150, 150)

    return slide

def generate_ppt():
    """PPT 생성 메인 함수"""

    print("[*] Git 입문 교재 PPT 생성 시작...")

    # 프레젠테이션 생성
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide_count = 1
    chapter_info = [
        (1, "Git 시작하기"),
        (2, "기본 작업"),
        (3, "협업하기"),
        (4, "문제 해결"),
        (5, "실전 상황"),
    ]

    for chapter_num, chapter_title in chapter_info:
        print(f"\n[*] Chapter {chapter_num}: {chapter_title} 처리 중...")

        html_file = f"chapter{chapter_num}.html"
        if not os.path.exists(html_file):
            print(f"   [!] {html_file} 파일을 찾을 수 없습니다.")
            continue

        chapter_colors = CHAPTER_COLORS[chapter_num]

        # 제목 슬라이드
        create_title_slide(prs, chapter_num, chapter_title, chapter_colors)
        print(f"   [OK] 제목 슬라이드 생성 (슬라이드 #{slide_count})")
        slide_count += 1

        # 섹션 슬라이드
        sections = extract_sections_from_html(html_file)
        for i, section in enumerate(sections):
            create_content_slide(prs, section, chapter_colors, slide_count)
            print(f"   [OK] '{section['title']}' 슬라이드 생성 (슬라이드 #{slide_count})")
            slide_count += 1

    # 저장
    output_path = "output/ppt/Git_입문_강의.pptx"
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)

    print(f"\n[OK] PPT 생성 완료!")
    print(f"[*] 저장 위치: {output_path}")
    print(f"[*] 총 슬라이드: {len(prs.slides)}")

if __name__ == "__main__":
    generate_ppt()
